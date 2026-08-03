import csv
import os

from langchain_community.document_loaders import PyPDFLoader

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx", ".txt", ".md", ".csv"}


def _extract_pdf_with_pypdf(file_path: str) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ValueError("Missing dependency for PDF extraction. Install `pypdf`.") from exc

    try:
        reader = PdfReader(file_path)
    except Exception as exc:
        raise ValueError(f"Could not read PDF file: {exc}") from exc

    if getattr(reader, "is_encrypted", False):
        try:
            reader.decrypt("")
        except Exception as exc:
            raise ValueError(
                "PDF is encrypted and cannot be processed without a password."
            ) from exc

    chunks = []
    for page in reader.pages:
        try:
            extracted = page.extract_text() or ""
        except Exception:
            extracted = ""
        if extracted.strip():
            chunks.append(extracted.strip())

    return "\n".join(chunks)


def _extract_pdf(file_path: str) -> str:
    loader_error = None
    try:
        loader = PyPDFLoader(file_path)
        pages = loader.load_and_split()
        text = " ".join(page.page_content for page in pages).strip()
        if text:
            return text
    except Exception as exc:
        loader_error = exc

    fallback_text = _extract_pdf_with_pypdf(file_path).strip()
    if fallback_text:
        return fallback_text

    if loader_error is not None:
        raise ValueError(f"Could not extract text from PDF. Parser error: {loader_error}")
    raise ValueError("Could not extract text from PDF.")


def _extract_docx(file_path: str) -> str:
    try:
        from docx import Document  # type: ignore
    except ImportError as exc:
        raise ValueError(
            "Missing dependency for DOCX extraction. Install `python-docx`."
        ) from exc

    document = Document(file_path)
    paragraph_text = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]

    table_text = []
    for table in document.tables:
        for row in table.rows:
            row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_values:
                table_text.append(" | ".join(row_values))

    return "\n".join(paragraph_text + table_text)


def _extract_xlsx(file_path: str) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore
    except ImportError as exc:
        raise ValueError(
            "Missing dependency for XLSX extraction. Install `openpyxl`."
        ) from exc

    workbook = load_workbook(filename=file_path, data_only=True, read_only=True)
    lines = []

    for sheet in workbook.worksheets:
        lines.append(f"[Sheet: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if values:
                lines.append(" | ".join(values))

    workbook.close()
    return "\n".join(lines)


def _extract_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError as exc:
        raise ValueError(
            "Missing dependency for PPTX extraction. Install `python-pptx`."
        ) from exc

    presentation = Presentation(file_path)
    lines = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"[Slide: {slide_index}]")
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    lines.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    row_values = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_values:
                        lines.append(" | ".join(row_values))

    return "\n".join(lines)


def _extract_csv(file_path: str) -> str:
    lines = []
    with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as file:
        reader = csv.reader(file)
        for row in reader:
            values = [value.strip() for value in row if value and value.strip()]
            if values:
                lines.append(" | ".join(values))
    return "\n".join(lines)


def _extract_plain_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as file:
        return file.read()


def extract_text_from_file(file_path: str) -> str:
    extension = os.path.splitext(file_path)[1].lower()

    if extension not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type '{extension}'. Supported: {supported}.")

    if extension == ".pdf":
        return _extract_pdf(file_path)
    if extension == ".docx":
        return _extract_docx(file_path)
    if extension == ".xlsx":
        return _extract_xlsx(file_path)
    if extension == ".pptx":
        return _extract_pptx(file_path)
    if extension == ".csv":
        return _extract_csv(file_path)
    return _extract_plain_text(file_path)
